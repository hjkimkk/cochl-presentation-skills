"""
Generates cochl_acousticshield_deck.pptx from the deck content.
Uses python-pptx. Run: python3 gen_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy, os

OUT = os.path.join(os.path.dirname(__file__), "cochl_acousticshield_deck.pptx")

# ── Color palette (Cochl PT template — Cochl_PT_template.pptx) ──────────
C_PRIMARY   = RGBColor(0x83, 0x2B, 0xFB)   # #832BFB — Cochl purple (from template)
C_ANALYTICS = RGBColor(0x2F, 0x80, 0xED)   # #2F80ED — secondary blue
C_LAVENDER  = RGBColor(0xEB, 0xDD, 0xFF)   # #EBDDFF — CTA subtitle
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG  = RGBColor(0xFF, 0xFF, 0xFF)   # content slide background
C_TEXT_DARK = RGBColor(0x33, 0x33, 0x33)   # body text on light slides
C_TEXT_SUB  = RGBColor(0x66, 0x66, 0x66)   # axis labels, footnotes
C_MUTED     = RGBColor(0x8B, 0x9B, 0xB4)
C_SUCCESS   = RGBColor(0x10, 0xB9, 0x81)
C_WARNING   = RGBColor(0xF5, 0x9E, 0x0B)
C_EMERGENCY = RGBColor(0xEB, 0x57, 0x57)
# Legacy aliases (kept for backward-compat with existing slide functions)
C_BG        = RGBColor(0x0B, 0x1F, 0x3A)   # dark bg for cover/break/CTA slides
C_CARD      = RGBColor(0x0F, 0x28, 0x47)
C_AI        = C_PRIMARY                      # AI purple = Cochl primary

# ── Slide dimensions (Cochl PT template: 10×5.625 in) ────────────────────
W = Inches(10.0)
H = Inches(5.625)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        text="", font_size=18, bold=False, color=C_WHITE,
        italic=False, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    try:
        run.font.name = "IBM Plex Sans"
    except Exception:
        run.font.name = "Calibri"
    return txBox


def rect(slide, left, top, width, height, fill_color=C_CARD, radius_pt=12):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()       # no line
    return shape


def accent_bar(slide, left, top, width=1.2):
    r = rect(slide, left, top, width, 0.04, fill_color=C_AI)
    return r


def label(slide, left, top, text):
    box(slide, left, top, 4, 0.22, text=text.upper(),
        font_size=10, color=C_MUTED)


def h2(slide, left, top, text, width=11.5):
    box(slide, left, top, width, 0.55, text=text,
        font_size=26, bold=True, color=C_WHITE)


def kpi_card(slide, left, top, w, h, value, lbl, sub=None):
    rect(slide, left, top, w, h, C_CARD)
    box(slide, left+0.15, top+0.12, w-0.3, 0.45,
        text=value, font_size=22, bold=True, color=C_AI)
    box(slide, left+0.15, top+0.55, w-0.3, 0.3,
        text=lbl, font_size=10, color=C_MUTED)
    if sub:
        box(slide, left+0.15, top+0.82, w-0.3, 0.18,
            text=sub, font_size=9, color=C_ANALYTICS)


def pillar_card(slide, left, top, w, h, num, heading, body):
    rect(slide, left, top, w, h, C_CARD)
    box(slide, left+0.15, top+0.15, w-0.3, 0.4,
        text=num, font_size=22, bold=True, color=C_AI)
    accent_bar(slide, left+0.15, top+0.55, 0.4)
    box(slide, left+0.15, top+0.65, w-0.3, 0.32,
        text=heading, font_size=12, bold=True, color=C_WHITE)
    box(slide, left+0.15, top+0.98, w-0.3, h-1.1,
        text=body, font_size=10, color=C_MUTED)


# ─────────────────────────────────────────────────────────
# SLIDE 1 · Cover
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
box(s, 0.7, 0.55, 10, 0.25, "DEFENSE R&D PROPOSAL · 방산혁신기업 지원사업 2026",
    font_size=10, color=C_MUTED)
box(s, 0.7, 0.85, 10, 1.1, "Cochl AcousticShield",
    font_size=48, bold=True, color=C_WHITE)
accent_bar(s, 0.7, 2.0, 1.2)
box(s, 0.7, 2.15, 10, 0.4,
    "광역 분산형 음향 AI 기반 드론 탐지 및 다중 위협 상황인식 체계",
    font_size=16, color=C_WHITE)
box(s, 0.7, 2.6, 9, 0.35,
    "When RF fails, when radar blinds, when cameras fog — acoustics listens.",
    font_size=14, italic=True, color=RGBColor(0xCC,0xCC,0xCC))
# meta row
metas = [("COMPANY","주식회사 코클"), ("DATE","2026.05.15"),
         ("DURATION","18개월 R&D"), ("FUNDING","정부지원금 20.0억원")]
for i,(lbl_t, val_t) in enumerate(metas):
    x = 0.7 + i * 3.0
    box(s, x, 3.2, 2.8, 0.2, lbl_t, font_size=9, color=C_MUTED)
    box(s, x, 3.42, 2.8, 0.3, val_t, font_size=13, bold=True, color=C_WHITE)
box(s, 0.7, 7.1, 1.5, 0.25, "Cochl", font_size=14, bold=True,
    color=RGBColor(0xCC,0xCC,0xCC))

# ─────────────────────────────────────────────────────────
# SLIDE 2 · Problem
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "CHALLENGE")
h2(s, 0.7, 0.6, "Modern Battlefield Has a Detection Gap")
cards = [
    ("📡","RF / Jamming — 구조적 무력화",
     "광섬유 FPV 드론은 RF 시그니처가 없다. 재머는 완전히 무력화됐다. 2024년 이후 전장의 표준."),
    ("🎯","Radar — 저고도·도심 사각",
     "낮은 RCS, 산악 클러터, 도심 다중반사로 소형 드론 식별 불가. 광역 저고도 사각 발생."),
    ("📷","EO / IR — 기상·광역 한계",
     "안개·야간 환경에서 무력화. 좁은 시야각으로 광역 탐색 구조적 불가."),
]
for i,(icon,heading,body) in enumerate(cards):
    x = 0.7 + i * 4.1
    rect(s, x, 1.5, 3.8, 3.8, C_CARD)
    box(s, x+0.2, 1.6, 3.4, 0.45, icon, font_size=24, color=C_WHITE)
    box(s, x+0.2, 2.12, 3.4, 0.35, heading, font_size=12, bold=True, color=C_WARNING)
    box(s, x+0.2, 2.52, 3.4, 1.6, body, font_size=11, color=C_MUTED)
rect(s, 0.7, 5.5, 11.9, 1.0, RGBColor(0x14,0x22,0x0A))
box(s, 0.9, 5.58, 11.5, 0.38,
    "2024: Russia deployed fiber-optic FPV drones — zero RF signature, fully jammer-immune. This is the new normal.",
    font_size=12, bold=True, color=C_WHITE)
box(s, 0.9, 5.98, 11.5, 0.3,
    "→ 우크라이나 전선에서 기존 RF/재머 체계의 구조적 한계 실전 확인",
    font_size=10, bold=True, color=C_WARNING)

# ─────────────────────────────────────────────────────────
# BREAK A · section-break
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
box(s, 0, 2.3, 13.33, 1.2, "02", font_size=96, bold=True,
    color=RGBColor(0x14,0x1E,0x2A), align=PP_ALIGN.CENTER)
box(s, 1, 2.8, 11.33, 1.0, "The Answer: Acoustics",
    font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
accent_bar(s, 6.2, 3.9, 0.9)
box(s, 1, 4.1, 11.33, 0.4,
    "드론이 비행하는 한, 음향 시그니처는 물리적으로 제거할 수 없다",
    font_size=14, color=C_MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────
# SLIDE 3 · Market
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "MARKET OPPORTUNITY")
h2(s, 0.7, 0.6, "A $20B Market at Its Inflection Point")
kpis = [
    ("$2.32B","드론 탐지 시장 (2029)","CAGR 28.7%"),
    ("$20.31B","Counter-UAS 시장 (2030)",None),
    ("~24,000","우크라이나 실전 배포 센서","Sky Fortress+Zvook"),
    ("15,000","NATO 추가 배치 (3세대)",None),
    ("~$10B","미 국방부 C-UAS 예산",None),
    ("$500M","FEMA C-UAS Grant","2025.12"),
]
for i,kpi in enumerate(kpis):
    row = i // 3; col = i % 3
    x = 0.7 + col * 4.1; y = 1.5 + row * 1.85
    kpi_card(s, x, y, 3.8, 1.65, kpi[0], kpi[1], kpi[2])
rect(s, 0.7, 5.35, 11.9, 0.65, RGBColor(0x0A,0x18,0x2A))
box(s, 1.1, 5.42, 11.3, 0.5,
    '"The U.S. homeland is no longer a sanctuary." — Pentagon, 2024',
    font_size=12, italic=True, color=C_WHITE)

# ─────────────────────────────────────────────────────────
# SLIDE 4 · Solution
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "SOLUTION")
h2(s, 0.7, 0.6, "Cochl AcousticShield: The Acoustic ISR Layer")
box(s, 0.7, 1.18, 11.9, 0.3,
    "음향 시그니처는 어떤 회피 기법으로도 제거할 수 없다 — 드론이 비행하는 한.",
    font_size=12, color=C_MUTED)
pillars = [
    ("01","저비용 분산 노드",
     "≤₩1.5M/노드 (1,000대 BOM 추정). 단일 마이크, 태양광+배터리, eSIM 글로벌 통신."),
    ("02","Edge AI + 확률적 위치추정",
     "TDOA 동기화 불요. 밀도 기반 추정으로 ±50m 위치 정확도. 100m 격자 배치."),
    ("03","6종 동시 위협 검출",
     "드론+총성+궤도차량+폭발음+침입음+음성 — 동일 노드·동일 모델."),
    ("04","대화형 AI Agent",
     "On-premise LLM(Gemma 3), ASR+TTS 음성 양방향, 자동 우선순위화 결심 지원."),
]
for i,(num,heading,body) in enumerate(pillars):
    x = 0.7 + i * 3.0
    pillar_card(s, x, 1.55, 2.75, 5.1, num, heading, body)

# ─────────────────────────────────────────────────────────
# SLIDE 5 · Architecture
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "ARCHITECTURE")
h2(s, 0.7, 0.6, "4계층 음향 ISR 아키텍처")
layers = [
    ("#2F80ED","Edge Layer — 분산 음향 센서 노드",
     "100m 간격 격자 분산, Edge AI 1차 분류, 이벤트 wake-up, 메타데이터(종류·신뢰도·타임스탬프)만 서버 전송"),
    ("#4B72F4","Server Layer — 융합·위치추정 서버",
     "노드 메타데이터 수신·정합, 밀도 기반 확률적 위치추정 (TDOA 동기화 불요, ±50m 이내)"),
    ("#5D5BF8","Fusion Layer — 통합 관제 + NATO SAPIENT API",
     "VMS 관제 화면(위협 지도·시계열·자동경보), NATO SAPIENT BSI Flex 335 v2 개방형 API"),
    ("#6B4EFF","Operator Layer — 대화형 AI Agent",
     "On-premise LLM(Gemma 3), ASR+TTS 음성 양방향, 이벤트 자동 요약·우선순위화, ≤15초 응답"),
]
colors_hex = ["#2F80ED","#4B72F4","#5D5BF8","#6B4EFF"]
colors_rgb = [RGBColor(0x2F,0x80,0xED), RGBColor(0x4B,0x72,0xF4),
              RGBColor(0x5D,0x5B,0xF8), RGBColor(0x6B,0x4E,0xFF)]
for i,(chex,heading,body) in enumerate(layers):
    y = 1.45 + i * 1.38
    rect(s, 0.7, y, 11.9, 1.18, C_CARD)
    # colored left border via thin rect
    rect(s, 0.7, y, 0.06, 1.18, colors_rgb[i])
    box(s, 0.95, y+0.1, 0.4, 0.4, str(i+1), font_size=18, bold=True,
        color=C_WHITE)
    box(s, 1.5, y+0.08, 10.8, 0.35, heading, font_size=13, bold=True,
        color=colors_rgb[i])
    box(s, 1.5, y+0.45, 10.8, 0.65, body, font_size=11, color=C_MUTED)
    if i < 3:
        box(s, 6.2, y+1.18, 0.9, 0.2, "▼", font_size=12, color=C_MUTED,
            align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────
# SLIDE 6 · Capabilities
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "CAPABILITIES")
h2(s, 0.7, 0.6, "One Network. Six Threat Classes.")
headers = ["위협 분류","탐지 대상","성능 목표"]
col_widths = [2.4, 6.8, 2.5]
row_data = [
    ("드론 존재 탐지","FPV급 ≥120m / Shahed-class ≥2,000m (SNR -5dB 환경)","≥90% 검출 정확도"),
    ("드론 모델명 분류 (10종)","DJI Mavic, Shahed, Wing Loong, Open-set detection","≥80% 분류 정확도"),
    ("총성","소화기, 기관총","≥85%"),
    ("궤도차량","탱크, 장갑차","≥85%"),
    ("침입음","유리 파괴, 펜스 절단","≥85%"),
    ("폭발음 / 포격","포병, 박격포 낙탄","≥85%"),
]
header_y = 1.42
for ci, (hdr, cw) in enumerate(zip(headers, col_widths)):
    x = 0.7 + sum(col_widths[:ci])
    box(s, x, header_y, cw, 0.3, hdr, font_size=9, color=C_MUTED)
for ri, row in enumerate(row_data):
    y = 1.8 + ri * 0.72
    fill = RGBColor(0x0F,0x28,0x47) if ri % 2 == 0 else RGBColor(0x0C,0x22,0x3C)
    rect(s, 0.7, y, 11.9, 0.62, fill)
    cols = list(row)
    for ci,(txt,cw) in enumerate(zip(cols, col_widths)):
        x = 0.7 + sum(col_widths[:ci])
        c = C_AI if ci == 0 else (C_SUCCESS if ci == 2 else C_MUTED)
        box(s, x+0.1, y+0.1, cw-0.2, 0.42, txt, font_size=11,
            bold=(ci==0), color=c)
rect(s, 0.7, 6.22, 0.06, 0.35, C_ANALYTICS)
box(s, 0.9, 6.22, 11.5, 0.3,
    "900+ 환경음 클래스 동시 학습으로 실외 오경보율 최소화. KOLAS 인증기관(KOTCA) 객관 검증.",
    font_size=10, color=C_MUTED)

# ─────────────────────────────────────────────────────────
# BREAK B · stat-spotlight
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
box(s, 0, 1.4, 13.33, 0.3, "DETECTION → DECISION · 강원 DMZ 광섬유 FPV 드론 시나리오",
    font_size=10, color=C_MUTED, align=PP_ALIGN.CENTER)
box(s, 0, 1.85, 13.33, 2.4, "17초",
    font_size=96, bold=True, color=C_AI, align=PP_ALIGN.CENTER)
box(s, 0, 4.35, 13.33, 0.55, "첫 검출에서 대응 결심까지",
    font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
box(s, 0, 5.0, 13.33, 0.3, "RF 0건 · 새벽 안개 · 03:14 · AI Agent 음성 보고 포함",
    font_size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────
# SLIDE 7 · Performance
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "PERFORMANCE TARGETS")
h2(s, 0.7, 0.6, "Quantified. KOLAS-Verified.")
perf_kpis = [
    ("≤500ms","Edge AI 추론 지연",None),
    ("≤15초","End-to-End 알람 시간",None),
    ("±50m","위치 추정 정확도","노드 간격 100m 기준"),
    ("≤10MB","Edge AI 모델 크기","INT8 양자화"),
    ("≤₩1.5M","노드 양산 단가","1,000대 BOM 추정"),
    ("≤15초","AI Agent 응답 시간",None),
    ("≥80%","AI Agent 시나리오 정확도",None),
    ("-20°~+50°C","동작 온도 범위","MIL-STD-810 준용"),
]
for i,kpi in enumerate(perf_kpis):
    row = i // 4; col = i % 4
    x = 0.7 + col * 3.0; y = 1.5 + row * 2.15
    kpi_card(s, x, y, 2.75, 1.9, kpi[0], kpi[1], kpi[2])
rect(s, 0.7, 5.9, 11.9, 0.45, RGBColor(0x0A,0x18,0x2A))
box(s, 0.9, 5.95, 11.5, 0.35,
    "전 항목 KOLAS 인증기관(KOTCA) 공인 성능평가로 객관 검증 · Phase C (M13~M18)",
    font_size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────
# SLIDE 8 · Scenarios
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "OPERATIONAL SCENARIOS")
h2(s, 0.7, 0.6, "17초 — 검출에서 대응 결심까지")
scenarios = [
    ("Scenario A","강원 DMZ — 광섬유 FPV (RF 0건)",
     ["03:14. 안개. 풍속 3m/s. RF 탐지 0건.",
      "노드 14대가 FPV 3기 순차 검출 (시속 90km)",
      '→ "동측 4.2km, FPV 3기. 예상 도달 2분30초. 즉시 대응 권고." — AI Agent',
      "첫 검출 → 대응 결심: 17초"]),
    ("Scenario B","휴전선 복합 침투 — 시계열 통합",
     ["02:15 보병 5명 외국어 무전 검출",
      "02:46 자폭 드론 동측 7.2km 검출",
      "03:18 궤도차량 후방 4km 검출",
      '→ "복합 침투 의심도 91%. 즉시 통합 대응." — AI Agent']),
    ("Scenario C","인천공항 — 동시 복합 위협",
     ["23:48 외곽 펜스 유리 파괴음 + 발걸음 검출",
      "동시 — 동측 1.8km 자폭 드론 접근",
      '→ "(1) 드론 차단 즉시. (2) Sector 7 동시 투입." — AI Agent']),
]
for i,(badge,heading,steps) in enumerate(scenarios):
    x = 0.7 + i * 4.1
    rect(s, x, 1.48, 3.8, 5.25, C_CARD)
    # badge
    rect(s, x+0.15, 1.62, 1.1, 0.28, C_AI)
    box(s, x+0.15, 1.62, 1.1, 0.28, badge, font_size=9, bold=True, color=C_WHITE)
    box(s, x+0.15, 1.97, 3.5, 0.45, heading, font_size=12, bold=True, color=C_WHITE)
    for j,step in enumerate(steps):
        box(s, x+0.15, 2.5+j*0.7, 3.5, 0.6, step, font_size=10, color=C_MUTED)

# ─────────────────────────────────────────────────────────
# SLIDE 9 · Competitive
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "COMPETITIVE LANDSCAPE")
h2(s, 0.7, 0.6, "차세대 vs 1세대 분산 음향 탐지")
comp_headers = ["기능","Sky Fortress\n우크라이나","Zvook\n우크라이나",
                "Squarehead G2+\n노르웨이","Cochl AcousticShield"]
comp_widths = [1.9, 2.2, 2.0, 2.35, 2.75]
for ci,(hdr,cw) in enumerate(zip(comp_headers, comp_widths)):
    x = 0.7 + sum(comp_widths[:ci])
    c = RGBColor(0xA0,0x8B,0xFF) if ci == 4 else C_MUTED
    box(s, x, 1.4, cw, 0.45, hdr, font_size=9, bold=(ci==0), color=c)
comp_rows = [
    ("위협 분류","드론 단독","드론 단독","드론 단독","드론 + 5종 군사 음향"),
    ("위치 추정","TDOA (clock sync)","TDOA","Phased array","밀도 기반 (동기화 불요)"),
    ("AI Agent","✗","✗","✗","On-premise LLM + 음성"),
    ("NATO SAPIENT","✗","✗","✗","BSI Flex 335 v2"),
    ("노드 단가","~$400-1,000","~$500","高 (array)","≤₩1.5M (~$1,100)"),
    ("실전 검증","✓ 95% 격추율","✓ 5km 탐지","DIU Finalist","R&D → KOLAS 검증"),
]
for ri, row in enumerate(comp_rows):
    y = 1.92 + ri * 0.72
    fill = RGBColor(0x0F,0x28,0x47) if ri % 2 == 0 else RGBColor(0x0C,0x22,0x3C)
    rect(s, 0.7, y, 11.2, 0.62, fill)
    # Cochl column highlight
    rect(s, 0.7+sum(comp_widths[:4]), y, comp_widths[4], 0.62,
         RGBColor(0x16,0x10,0x38))
    for ci,(txt,cw) in enumerate(zip(row, comp_widths)):
        x = 0.7 + sum(comp_widths[:ci])
        if ci == 0: c = C_WHITE
        elif txt == "✗": c = C_EMERGENCY
        elif txt.startswith("✓"): c = C_SUCCESS
        elif ci == 4: c = C_SUCCESS
        else: c = C_MUTED
        box(s, x+0.08, y+0.12, cw-0.16, 0.38, txt, font_size=10, color=c)
rect(s, 0.7, 6.24, 0.06, 0.3, C_ANALYTICS)
box(s, 0.9, 6.24, 11, 0.28,
    "비교 기준: 공개 데이터시트, 학술 논문, 현장 보고서 (2025~2026). Cochl KOLAS 검증 후 정량 목표 달성 예정.",
    font_size=9, color=C_MUTED)

# ─────────────────────────────────────────────────────────
# BREAK C · quote-break
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
box(s, 0, 0.5, 13.33, 2.0, "“", font_size=96, bold=True,
    color=RGBColor(0x18,0x22,0x30), align=PP_ALIGN.LEFT)
box(s, 1.5, 2.2, 10.33, 1.6,
    '"The U.S. homeland is no longer a sanctuary."',
    font_size=30, italic=True, color=C_WHITE, align=PP_ALIGN.CENTER)
box(s, 1.5, 3.9, 10.33, 0.4, "— Pentagon, 2024",
    font_size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
accent_bar(s, 6.1, 4.5, 1.1)

# ─────────────────────────────────────────────────────────
# SLIDE 10 · Technology
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "TECHNOLOGY & CREDENTIALS")
h2(s, 0.7, 0.6, "9년의 사운드 AI. 이제 전장에서 검증한다.")
cols_data = [
    ("Cochl IP Stack", C_ANALYTICS, [
        "Cochl.Sense Foundation Model — TRL 7, production 검증",
        "900+ 클래스 환경음 DB (9년 누적 구축)",
        "Transformer + Knowledge Distillation 경량화 파이프라인",
        "CNN-Conformer 등 학계 SOTA 비교 검증 병행",
        "MLOps 풀스택 (수집→레이블→학습→배포→모니터링)",
        "국제특허 12건 등록 · 42건 출원",
    ]),
    ("Proven Track Record", C_AI, [
        "IEEE DCASE 2년 연속 우승",
        "NVIDIA Inception Award 2018",
        "미 해군(US Navy) 기지 보안 배포",
        "미 공군 Nellis 기지 (Ghost Robotics 협업)",
        "USSOCOM 데모 — 2026년 3월",
        "LIG D&A · 한화시스템 잠수함 음향 분석 협업",
        "방산혁신기업100 AI분야 선정 (2024.09)",
    ]),
]
for i,(heading,hcolor,items) in enumerate(cols_data):
    x = 0.7 + i * 6.15
    rect(s, x, 1.45, 5.8, 5.5, C_CARD)
    box(s, x+0.2, 1.58, 5.4, 0.35, heading, font_size=12,
        bold=True, color=hcolor)
    for j,item in enumerate(items):
        box(s, x+0.2, 2.05+j*0.68, 5.4, 0.58,
            "▸  " + item, font_size=11, color=C_MUTED)

# ─────────────────────────────────────────────────────────
# SLIDE 11 · Roadmap
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "R&D ROADMAP")
h2(s, 0.7, 0.6, "18개월 단계별 로드맵")
# timeline track
rect(s, 0.7, 1.45, 11.9, 0.05, C_AI)
phases = [
    ("Phase 1","M1-M2\nSep-Oct '26","Kickoff",
     "과제 착수, DB 구축,\nAI 모델 v1 시작","₩2.2억",
     RGBColor(0x4B,0x68,0xFF), False),
    ("Phase 2","M3-M4\nNov-Dec '26","Alpha",
     "AI 모델 v1 alpha,\n시제 2대 제작","₩2.3억",
     RGBColor(0x5B,0x78,0xFF), False),
    ("Phase 3","M5-M7\nJan-Mar '27","Integration",
     "시제 6대 통합,\n모델 v2 착수 ▶ 중간점검","₩3.5억",
     C_WARNING, True),
    ("Phase 4","M8-M10\nApr-Jun '27","Field Test 1",
     "수도권 야외 시험,\nAI Agent v1, SAPIENT v1","₩3.8억",
     RGBColor(0x6D,0x56,0xFA), False),
    ("Phase 5","M11-M13\nJul-Sep '27","Full Integration",
     "전체 통합,\nAI Agent v2(음성)","₩3.8억",
     RGBColor(0x83,0x2B,0xFB), False),
    ("Phase 6","M14-M16\nOct-Dec '27","KOLAS",
     "SAPIENT 검증,\nKOLAS 시험 ▶ 성과평가","₩2.2억",
     C_WARNING, True),
    ("Phase 7","M17-M18\nJan-Feb '28","Certification",
     "KOLAS 인증 통과,\n특허·논문, 최종 평가","₩2.2억",
     C_SUCCESS, False),
]
for i,ph in enumerate(phases):
    pid,period,lbl,desc,budget,color,chk = ph[0],ph[1],ph[2],ph[3],ph[4],ph[5],(len(ph)>6 and ph[6])
    x = 0.7 + i * 1.72
    # dot
    dot = slide.shapes if False else None
    rect(s, x+0.75, 1.3, 0.2, 0.2, color)
    box(s, x, 1.58, 1.6, 0.28, pid, font_size=9, bold=True, color=color)
    box(s, x, 1.88, 1.6, 0.35, period, font_size=8, color=C_MUTED)
    box(s, x, 2.28, 1.6, 0.3, lbl, font_size=10, bold=True, color=C_WHITE)
    box(s, x, 2.62, 1.6, 1.5, desc, font_size=9, color=C_MUTED)
    rect(s, x, 4.22, 1.5, 0.3, C_CARD)
    box(s, x, 4.22, 1.5, 0.3, budget, font_size=9, bold=True, color=color)
rect(s, 0.7, 4.72, 11.9, 0.5, C_CARD)
box(s, 4.0, 4.77, 5.33, 0.38,
    "총 예산 ₩26.7억 · 18개월 단일 단계",
    font_size=13, bold=True, color=C_AI, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────
# SLIDE 12 · Budget
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
label(s, 0.7, 0.4, "BUDGET")
h2(s, 0.7, 0.6, "총 예산 ₩26.7억 R&D 투자")
# Donut (described in text since we can't easily draw SVG in pptx)
rect(s, 0.7, 1.5, 4.2, 4.5, C_CARD)
box(s, 0.9, 1.65, 3.8, 0.35, "예산 구조", font_size=11, bold=True, color=C_MUTED)
box(s, 1.0, 2.1, 3.6, 0.55, "75%", font_size=40, bold=True, color=C_AI,
    align=PP_ALIGN.CENTER)
box(s, 1.0, 2.7, 3.6, 0.3, "정부지원금", font_size=12, color=C_MUTED,
    align=PP_ALIGN.CENTER)
box(s, 1.0, 3.1, 3.6, 0.55, "₩20.0억", font_size=20, bold=True, color=C_WHITE,
    align=PP_ALIGN.CENTER)
rect(s, 1.0, 3.8, 0.18, 0.18, C_AI)
box(s, 1.25, 3.8, 2.8, 0.22, "정부지원금 75% — ₩20.0억", font_size=10, color=C_MUTED)
rect(s, 1.0, 4.08, 0.18, 0.18, C_ANALYTICS)
box(s, 1.25, 4.08, 2.8, 0.22, "기업부담금 25% — ₩6.7억", font_size=10, color=C_MUTED)
# Bars
rect(s, 5.2, 1.5, 7.2, 4.5, C_CARD)
box(s, 5.4, 1.65, 6.8, 0.3, "예산 배분", font_size=11, bold=True, color=C_MUTED)
bar_rows = [
    ("AI 연구 및 모델 개발", 35),
    ("SW (융합서버, VMS, SAPIENT API)", 20),
    ("HW 통합 (시제 6대 제작)", 20),
    ("야외 시험 및 KOLAS 인증", 15),
    ("AI Agent 및 LLM 통합", 10),
]
for j,(blbl,pct) in enumerate(bar_rows):
    y = 2.1 + j * 0.72
    box(s, 5.4, y, 4.5, 0.28, blbl, font_size=10, color=C_MUTED)
    box(s, 9.95, y, 0.55, 0.28, f"{pct}%", font_size=10, bold=True, color=C_WHITE)
    rect(s, 5.4, y+0.32, 4.7, 0.22, RGBColor(0x14,0x22,0x32))
    fill_w = 4.7 * pct / 100
    rect(s, 5.4, y+0.32, fill_w, 0.22, C_AI)
rect(s, 0.7, 6.15, 0.06, 0.42, C_SUCCESS)
box(s, 0.9, 6.15, 11.5, 0.38,
    "양산 단계: 1,000노드 ≤₩15억 — 단일 레이더 수십~수백억원 대비 비대칭 방어 경제성 확보",
    font_size=11, color=RGBColor(0xCC,0xFF,0xCC))

# ─────────────────────────────────────────────────────────
# SLIDE 13 · CTA
# ─────────────────────────────────────────────────────────
s = add_slide(); bg(s, C_BG)
box(s, 0.7, 0.45, 11.9, 0.65,
    "한국 최초 광역 음향 ISR 레이어 구축",
    font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
box(s, 1.5, 1.15, 10.33, 0.3,
    "코클이 9년간 쌓은 사운드 AI를 이제 전장에서 검증합니다.",
    font_size=12, color=C_MUTED, align=PP_ALIGN.CENTER)
outcomes = [
    ("01","국가 안보","광섬유 드론 시대, 한국 최초 광역 분산 음향 ISR 자산 확보"),
    ("02","글로벌 수출","NATO SAPIENT+eSIM → NATO, USCENTCOM, 터키, 동남아 직접 수출"),
    ("03","시장 선점","$20B C-UAS 시장, 음향 AI 기반 방어 가능한 글로벌 진입 경로"),
]
for i,(num,heading,body) in enumerate(outcomes):
    x = 0.7 + i * 4.0
    rect(s, x, 1.6, 3.7, 2.5, RGBColor(0x14,0x20,0x38))
    box(s, x+0.2, 1.72, 3.3, 0.55, num, font_size=28, bold=True, color=C_AI)
    box(s, x+0.2, 2.35, 3.3, 0.38, heading, font_size=13, bold=True, color=C_WHITE)
    box(s, x+0.2, 2.78, 3.3, 1.0, body, font_size=11, color=C_MUTED)
rect(s, 1.5, 4.3, 10.33, 0.65, RGBColor(0x12,0x10,0x2C))
box(s, 1.5, 4.35, 10.33, 0.55,
    "코클이 9년간 쌓은 사운드 AI를 이제 전장에서 검증합니다.",
    font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
pipeline_steps = ["과제 선정","18개월 R&D","KOLAS 검증","신속시범획득","글로벌 수출"]
total_w = 11.9; step_w = 1.8; gap = 0.55
for i,step in enumerate(pipeline_steps):
    x = 0.7 + i * (step_w + gap)
    c = C_AI if i == len(pipeline_steps)-1 else C_CARD
    rect(s, x, 5.15, step_w, 0.45, c)
    tc = C_WHITE
    box(s, x, 5.15, step_w, 0.45, step, font_size=10,
        bold=(i==len(pipeline_steps)-1), color=tc, align=PP_ALIGN.CENTER)
    if i < len(pipeline_steps)-1:
        box(s, x+step_w, 5.24, gap, 0.28, "→", font_size=14,
            color=C_MUTED, align=PP_ALIGN.CENTER)
box(s, 0.7, 6.98, 11.9, 0.3, "contact@cochlear.ai · cochl.ai",
    font_size=10, color=RGBColor(0x66,0x66,0x66), align=PP_ALIGN.CENTER)
box(s, 0.7, 7.1, 1.5, 0.25, "Cochl", font_size=13, bold=True,
    color=RGBColor(0xAA,0xAA,0xAA))

# ─────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✓ Saved: {OUT}")
print(f"  Slides: {len(prs.slides)}")
